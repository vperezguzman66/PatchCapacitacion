import argparse
from dataclasses import dataclass
import unicodedata
from typing import Dict, List, Set

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ── Colores ──────────────────────────────────────────────
AZUL_OSCURO = RGBColor(0x1F, 0x4E, 0x79)
VERDE = RGBColor(0x17, 0x86, 0x4E)
AMARILLO = RGBColor(0xF1, 0xC4, 0x0F)
NARANJO = RGBColor(0xE6, 0x7E, 0x22)
GRIS_CLARO = RGBColor(0xF2, 0xF2, 0xF2)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_TEXTO = RGBColor(0x44, 0x44, 0x44)


@dataclass
class Recomendacion:
    tecnologia: str
    categoria: str
    familia: str
    prioridad: float
    nivel: str
    apps: int
    relacionadas: str
    afinidad: str
    accion: str


CATALOGO_CURSOS = [
    {
        "plataforma": "Udemy",
        "claves": ["python"],
        "tecnologia": "Python",
        "curso": "100 Days of Code™: The Complete Python Pro Bootcamp",
        "duracion": "",
    },
    {
        "plataforma": "Udemy",
        "claves": ["python"],
        "tecnologia": "Python",
        "curso": "Automate the Boring Stuff with Python Programming",
        "duracion": "",
    },
    {
        "plataforma": "Udemy",
        "claves": ["java"],
        "tecnologia": "Java",
        "curso": "Java Masterclass 2025: 130+ Hours of Expert Lessons",
        "duracion": "130+ horas",
    },
    {
        "plataforma": "Udemy",
        "claves": ["java"],
        "tecnologia": "Java",
        "curso": "Modern Java: Mastering Features from Java 8 to Java 25",
        "duracion": "",
    },
    {
        "plataforma": "Udemy",
        "claves": ["spring boot", "spring framework", "spring mvc", "spring cloud"],
        "tecnologia": "Spring Boot",
        "curso": "Master Spring Boot 3 & Spring Framework 6 with Java",
        "duracion": "",
    },
    {
        "plataforma": "Udemy",
        "claves": ["spring boot", "microservices"],
        "tecnologia": "Spring Boot",
        "curso": "Master Microservices with Spring Boot and Spring Cloud",
        "duracion": "",
    },
    {
        "plataforma": "Udemy",
        "claves": ["docker"],
        "tecnologia": "Docker",
        "curso": "Docker for the Absolute Beginner - Hands On - DevOps",
        "duracion": "",
    },
    {
        "plataforma": "Udemy",
        "claves": ["docker"],
        "tecnologia": "Docker",
        "curso": "Docker & Kubernetes: The Practical Guide",
        "duracion": "",
    },
    {
        "plataforma": "Udemy",
        "claves": ["kubernetes"],
        "tecnologia": "Kubernetes",
        "curso": "Certified Kubernetes Administrator (CKA) with Practice Tests",
        "duracion": "",
    },
    {
        "plataforma": "Udemy",
        "claves": ["kubernetes"],
        "tecnologia": "Kubernetes",
        "curso": "Kubernetes Certified Application Developer (CKAD) with Tests",
        "duracion": "",
    },
    {
        "plataforma": "Udemy",
        "claves": ["aws", "amazon web services", "amazon aws"],
        "tecnologia": "AWS",
        "curso": "Ultimate AWS Certified Cloud Practitioner CLF-C02 2026",
        "duracion": "",
    },
    {
        "plataforma": "Udemy",
        "claves": ["aws", "amazon web services", "amazon aws"],
        "tecnologia": "AWS",
        "curso": "AWS From Zero to Hero - The Complete Guide",
        "duracion": "",
    },
    {
        "plataforma": "Udemy",
        "claves": ["azure", "az-900", "microsoft azure"],
        "tecnologia": "Azure",
        "curso": "Master Microsoft Azure Fundamentals: AZ-900 Exam Prep 2026",
        "duracion": "",
    },
    {
        "plataforma": "Udemy",
        "claves": ["azure", "az-900", "microsoft azure"],
        "tecnologia": "Azure",
        "curso": "Microsoft Azure Fundamentals: AZ-900 Full Course & Exams",
        "duracion": "",
    },
    {
        "plataforma": "Udemy",
        "claves": ["sql", "mysql", "postgresql", "database"],
        "tecnologia": "SQL",
        "curso": "The Complete SQL Bootcamp: Go from Zero to Hero",
        "duracion": "30 horas",
    },
    {
        "plataforma": "Udemy",
        "claves": ["sql", "mysql", "postgresql", "database"],
        "tecnologia": "SQL",
        "curso": "SQL and PostgreSQL: The Complete Developer's Guide",
        "duracion": "",
    },
    {
        "plataforma": "Udemy",
        "claves": ["mongodb", "nosql"],
        "tecnologia": "MongoDB",
        "curso": "MongoDB - The Complete Developer's Guide",
        "duracion": "",
    },
    {
        "plataforma": "Udemy",
        "claves": ["mongodb", "nosql"],
        "tecnologia": "MongoDB",
        "curso": "The Complete MongoDB Course",
        "duracion": "",
    },
]


def semaforo_afinidad(conocidas_cat: int, total_cat: int) -> str:
    ratio = conocidas_cat / max(total_cat, 1)
    if conocidas_cat >= 2 or ratio >= 0.34:
        return "🟢 Alta"
    if conocidas_cat == 1 or ratio >= 0.12:
        return "🟡 Media"
    return "🔴 Baja"


def normalizar(texto: str) -> str:
    base = unicodedata.normalize("NFKD", str(texto))
    return "".join(c for c in base if not unicodedata.combining(c)).lower().strip()


def buscar_curso(tecnologia: str, categoria: str) -> str | None:
    texto = normalizar(tecnologia)
    cat = normalizar(categoria)
    for item in CATALOGO_CURSOS:
        claves = item["claves"]
        if any(clave in texto or clave in cat for clave in claves):
            duracion = f" · {item['duracion']}" if item.get("duracion") else ""
            return f"{item['plataforma']} · {item['tecnologia']} → {item['curso']}{duracion}"
    return None


def familia_tecnologia(tecnologia: str, categoria: str = "") -> str:
    texto = normalizar(f"{categoria} {tecnologia}")
    if any(k in texto for k in ["seguridad", "security", "ciber"]):
        return "Seguridad"
    if any(k in texto for k in ["infra", "devops", "kubernetes", "docker", "container", "contenedor"]):
        return "DevOps/Infra"
    if any(k in texto for k in ["cloud", "azure", "aws", "gcp", "almacenamiento", "storage", "nube"]):
        return "Cloud"
    if any(k in texto for k in ["datos", "base de datos", "database", "sql", "mysql", "postgres", "mongodb", "nosql", "oracle"]):
        return "Datos"
    if any(k in texto for k in ["api", "integr", "microservice", "microserv", "rest", "soap"]):
        return "Integración/API"
    if any(k in texto for k in ["framework", "spring", "hibernate", "jpa", "django", "flask", "react", "angular", "vue", "testing", "test"]):
        return "Frameworks"
    if any(k in texto for k in ["lenguaje", "programacion", "programación", "java", "python", "javascript", "typescript", "c#", "go", "kotlin"]):
        return "Lenguajes"
    return "Otros"


FAMILIAS_FOCUS = {
    "Lenguajes": {"Lenguajes", "Frameworks", "Datos", "Integración/API"},
    "Frameworks": {"Frameworks", "Lenguajes", "Datos", "Integración/API"},
    "Datos": {"Datos", "Lenguajes", "Frameworks", "Integración/API"},
    "Integración/API": {"Integración/API", "Frameworks", "Lenguajes", "Datos"},
    "Cloud": {"Cloud", "Frameworks", "Datos", "Integración/API"},
    "DevOps/Infra": {"DevOps/Infra", "Cloud"},
    "Seguridad": {"Seguridad"},
    "Otros": {"Otros"},
}


def familias_de_enfoque(familias_conocidas: Set[str]) -> Set[str]:
    foco: Set[str] = set()
    for familia in familias_conocidas:
        foco.update(FAMILIAS_FOCUS.get(familia, {familia}))
    if "Seguridad" not in familias_conocidas:
        foco.discard("Seguridad")
    if "DevOps/Infra" not in familias_conocidas:
        foco.discard("DevOps/Infra")
    return foco or set(familias_conocidas)


def base_conocida_texto(conocidas: Set, tech_categoria: Dict) -> str:
    tecnologias = sorted([str(t).strip() for t in conocidas], key=str.lower)
    if not tecnologias:
        return "Sin base conocida relevante"
    detalle = []
    for tech in tecnologias[:8]:
        familia = familia_tecnologia(tech, str(tech_categoria.get(tech, "")).strip())
        detalle.append(f"{tech} [{familia}]")
    return "\n".join([f"• {item}" for item in detalle])


def add_rect(slide, l, t, w, h, fill_rgb):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def add_text(
    slide,
    text,
    l,
    t,
    w,
    h,
    size=10,
    bold=False,
    color=RGBColor(0, 0, 0),
    align=PP_ALIGN.LEFT,
):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def cargar_datos(archivo: str):
    raw_matriz = pd.read_excel(archivo, sheet_name="Matriz", header=None)
    categorias_row = raw_matriz.iloc[1, 2:].ffill().tolist()
    tecnologias_nombres = raw_matriz.iloc[2, 2:].tolist()
    tech_categoria = {t: c for t, c in zip(tecnologias_nombres, categorias_row)}

    matriz = pd.read_excel(archivo, sheet_name="Matriz", header=2, index_col=1)
    matriz = matriz.drop(columns=[matriz.columns[0]], errors="ignore").dropna(how="all").fillna(0)

    personas = pd.read_excel(archivo, sheet_name="MtzPerson", header=2, index_col=1)
    personas = personas.drop(columns=[personas.columns[0]], errors="ignore").dropna(how="all").fillna(0)

    tecnologias = matriz.columns.intersection(personas.columns)
    matriz = matriz[tecnologias]
    personas = personas[tecnologias]

    techs_relevantes = set(tecnologias[(matriz == 1).any(axis=0)])

    return matriz, personas, tecnologias, techs_relevantes, tech_categoria


def accion_por_categoria(categoria: str) -> str:
    c = str(categoria).lower()
    if "lenguaje" in c:
        return "Curso avanzado + mini proyecto"
    if "datos" in c or "base de datos" in c:
        return "Laboratorio de datos aplicado"
    if "seguridad" in c:
        return "Secure coding + hardening"
    if "integr" in c or "api" in c:
        return "Taller de integración end-to-end"
    if "cloud" in c or "infra" in c or "plataforma" in c:
        return "Ruta hands-on con sandbox"
    return "Curso guiado + mentoring"


def proponer_para_persona(
    persona_row: pd.Series,
    tecnologias: pd.Index,
    techs_relevantes: Set,
    tech_categoria: Dict,
    demanda_apps: Dict,
    familias_objetivo: Set[str],
) -> List[Recomendacion]:
    conocidas = set(tecnologias[persona_row == 1]) & techs_relevantes
    gaps = list(techs_relevantes - conocidas)

    if not gaps:
        return []

    max_demanda = max(demanda_apps.values()) if demanda_apps else 1

    conocidas_por_cat = {}
    total_por_cat = {}

    for t in techs_relevantes:
        cat = str(tech_categoria.get(t, "Otros")).strip()
        total_por_cat[cat] = total_por_cat.get(cat, 0) + 1

    for t in conocidas:
        cat = str(tech_categoria.get(t, "Otros")).strip()
        conocidas_por_cat.setdefault(cat, set()).add(t)

    recs = []
    for t in gaps:
        cat = str(tech_categoria.get(t, "Otros")).strip()
        familia = familia_tecnologia(str(t), cat)
        if familia not in familias_objetivo:
            continue
        apps = int(demanda_apps.get(t, 0))
        score_demanda = apps / max_demanda
        conocidas_cat = len(conocidas_por_cat.get(cat, set()))
        total_cat = max(total_por_cat.get(cat, 1), 1)
        score_afinidad = conocidas_cat / total_cat
        prioridad = (0.6 * score_demanda) + (0.4 * score_afinidad)

        if conocidas_cat >= 2:
            nivel = "Intermedio"
        elif conocidas_cat == 1:
            nivel = "Básico-Intermedio"
        else:
            nivel = "Fundamentos"

        afinidad = semaforo_afinidad(conocidas_cat, total_cat)

        rel = sorted([str(x).strip() for x in conocidas_por_cat.get(cat, set())])
        relacionadas = ", ".join(rel[:3]) if rel else "Sin base directa"

        recs.append(
            Recomendacion(
                tecnologia=str(t).strip(),
                categoria=cat,
                familia=familia,
                prioridad=round(prioridad, 4),
                nivel=nivel,
                apps=apps,
                relacionadas=relacionadas,
                afinidad=afinidad,
                accion=accion_por_categoria(cat),
            )
        )

    recs.sort(key=lambda r: (-r.prioridad, -r.apps, r.tecnologia.lower()))
    return recs


def resumen_fase_texto(recs_fase: List[Recomendacion], limite: int = 4) -> str:
    if not recs_fase:
        return "Sin recomendaciones en esta fase"
    top = recs_fase[:limite]
    return "\n".join(
        (
            f"• {r.tecnologia} ({r.categoria}) [{r.familia}] | {r.afinidad} | Nivel: {r.nivel} | Apps: {r.apps}\n"
            f"  ↳ Intersección base: {r.relacionadas}"
        )
        for r in top
    )


def cursos_sugeridos_texto(conocidas: Set, tech_categoria: Dict, demanda_apps: Dict, limite: int = 3) -> str:
    if not conocidas:
        return "Perfeccionamiento sugerido: sin base conocida relevante para asociar cursos."

    tecnologias_ordenadas = sorted(
        [str(t).strip() for t in conocidas],
        key=lambda t: (-int(demanda_apps.get(t, 0)), t.lower()),
    )

    sugerencias = []
    vistos = set()
    for tecnologia in tecnologias_ordenadas:
        categoria = str(tech_categoria.get(tecnologia, "Otros")).strip()
        curso = buscar_curso(tecnologia, categoria)
        if curso and curso not in vistos:
            sugerencias.append(curso)
            vistos.add(curso)
        if len(sugerencias) >= limite:
            break

    if not sugerencias:
        return "Perfeccionamiento sugerido: revisar catálogo interno o rutas de formación equivalentes."

    return "\n".join([f"• {item}" for item in sugerencias])


def generar_presentacion(archivo_entrada: str, archivo_salida: str):
    matriz, personas, tecnologias, techs_relevantes, tech_categoria = cargar_datos(archivo_entrada)

    demanda_apps = {
        t: int((pd.to_numeric(matriz[t], errors="coerce").fillna(0) == 1).sum())
        for t in tecnologias
        if t in techs_relevantes
    }

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Portada
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.33, 7.5, AZUL_OSCURO)
    add_text(s, "Propuesta de Capacitación", 1, 2.0, 11, 1.2, size=40, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    add_text(
        s,
        "Plan personalizado por persona (Quick Wins, Consolidación y Expansión)",
        1,
        3.3,
        11,
        0.8,
        size=17,
        color=RGBColor(0xBD, 0xD7, 0xEE),
        align=PP_ALIGN.CENTER,
    )
    add_text(s, "Banco Santander Chile · 2026", 1, 4.2, 11, 0.35, size=14, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.CENTER)

    for persona_nombre, persona_row in personas.iterrows():
        conocidas = set(tecnologias[persona_row == 1]) & techs_relevantes
        familias_conocidas = {
            familia_tecnologia(str(t).strip(), str(tech_categoria.get(t, "")).strip())
            for t in conocidas
        }
        familias_objetivo = familias_de_enfoque(familias_conocidas)
        familias_conocidas_txt = " | ".join(sorted(familias_conocidas)) if familias_conocidas else "Sin familias claras"
        familias_objetivo_txt = " | ".join(sorted(familias_objetivo)) if familias_objetivo else "Sin familias objetivo"

        recs = proponer_para_persona(
            persona_row=persona_row,
            tecnologias=tecnologias,
            techs_relevantes=techs_relevantes,
            tech_categoria=tech_categoria,
            demanda_apps=demanda_apps,
            familias_objetivo=familias_objetivo,
        )

        f1 = recs[:6]
        f2 = recs[6:12]
        f3 = recs[12:20]

        gaps = set(techs_relevantes) - conocidas
        conocidas_ordenadas = sorted([str(t).strip() for t in conocidas], key=str.lower)

        conocidas_por_cat = {}
        for t in conocidas:
            cat = str(tech_categoria.get(t, "Otros")).strip()
            conocidas_por_cat[cat] = conocidas_por_cat.get(cat, 0) + 1
        top_cat = sorted(conocidas_por_cat.items(), key=lambda x: (-x[1], x[0].lower()))
        top_cat_txt = " | ".join([f"{c} ({n})" for c, n in top_cat[:4]]) if top_cat else "Sin categorías destacadas"

        s = prs.slides.add_slide(blank)
        add_rect(s, 0, 0, 13.33, 7.5, GRIS_CLARO)
        add_rect(s, 0, 0, 13.33, 1.0, AZUL_OSCURO)

        add_text(s, str(persona_nombre), 0.3, 0.08, 8.8, 0.55, size=22, bold=True, color=BLANCO)
        add_text(
            s,
            f"Conocidas: {len(conocidas)} | Brechas enfocadas: {len(gaps)} | Recomendadas: {len(recs)}",
            0.3,
            0.62,
            10,
            0.24,
            size=10,
            color=RGBColor(0xBD, 0xD7, 0xEE),
        )
        add_text(
            s,
            f"Base/familia: {familias_conocidas_txt}",
            0.3,
            0.85,
            10.5,
            0.18,
            size=9,
            bold=True,
            color=RGBColor(0xE6, 0xF4, 0xFA),
        )
        add_text(
            s,
            f"Alcance de la propuesta: {familias_objetivo_txt}",
            0.3,
            1.03,
            10.5,
            0.18,
            size=8,
            bold=True,
            color=RGBColor(0xD7, 0xEE, 0xF7),
        )

        # Bloque principal: perfeccionamiento sugerido y base conocida explícita
        add_rect(s, 0.2, 1.35, 12.95, 5.95, RGBColor(0xF7, 0xFB, 0xFF))
        add_rect(s, 0.28, 1.44, 12.79, 0.56, RGBColor(0x1D, 0x72, 0x8A))
        foco = recs[0] if recs else None
        if foco:
            foco_txt = (
                f"Foco inmediato sugerido: {foco.tecnologia} ({foco.categoria}) | "
                f"{foco.afinidad} | Nivel: {foco.nivel} | Apps: {foco.apps} | Intersección: {foco.relacionadas} | Acción: {foco.accion}"
            )
        else:
            foco_txt = "Sin brechas detectadas en tecnologías relevantes del banco."

        cursos_txt = cursos_sugeridos_texto(conocidas, tech_categoria, demanda_apps, limite=3)

        add_text(
            s,
            "PERFECCIONAMIENTO SUGERIDO",
            0.35,
            1.58,
            12.6,
            0.16,
            size=9,
            bold=True,
            color=BLANCO,
            align=PP_ALIGN.LEFT,
        )
        add_rect(s, 0.34, 2.06, 12.65, 1.62, BLANCO)
        add_text(
            s,
            cursos_txt,
            0.35,
            2.14,
            12.6,
            1.42,
            size=8,
            bold=False,
            color=RGBColor(0x1F, 0x4E, 0x79),
        )
        add_rect(s, 0.34, 3.78, 12.65, 1.30, RGBColor(0xE8, 0xF4, 0xF8))
        add_text(
            s,
            "BASE CONOCIDA EXPLÍCITA",
            0.35,
            3.86,
            12.6,
            0.16,
            size=8,
            bold=True,
            color=RGBColor(0x1D, 0x4E, 0x2E),
        )
        add_text(
            s,
            base_conocida_texto(conocidas, tech_categoria),
            0.35,
            4.06,
            12.6,
            0.92,
            size=7,
            bold=False,
            color=RGBColor(0x1D, 0x4E, 0x2E),
        )
        add_rect(s, 0.34, 5.16, 12.65, 1.00, RGBColor(0xF4, 0xF6, 0xF7))
        add_text(s, "FOCO INMEDIATO", 0.35, 5.24, 12.6, 0.14, size=8, bold=True, color=RGBColor(0x5B, 0x6B, 0x73))
        add_text(
            s,
            foco_txt,
            0.35,
            5.40,
            12.6,
            0.66,
            size=7,
            bold=False,
            color=RGBColor(0x5B, 0x6B, 0x73),
        )

    output = archivo_salida
    try:
        prs.save(output)
    except PermissionError:
        output = "PropuestaCapacitacion_actualizada.pptx"
        prs.save(output)

    print(f"Presentación generada: {output} ({len(personas)} personas)")


def main():
    parser = argparse.ArgumentParser(description="Genera propuesta de capacitación personalizada en formato PPT")
    parser.add_argument("--input", default="AplicacionesBanco.xlsx", help="Excel de entrada")
    parser.add_argument("--output", default="PropuestaCapacitacion.pptx", help="PPT de salida")
    args = parser.parse_args()
    generar_presentacion(args.input, args.output)


if __name__ == "__main__":
    main()
