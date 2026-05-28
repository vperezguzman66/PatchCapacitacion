import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

archivo = "AplicacionesBanco.xlsx"

# ── Colores ──────────────────────────────────────────────
AZUL_OSCURO = RGBColor(0x1F, 0x4E, 0x79)
AZUL_MEDIO  = RGBColor(0x27, 0x72, 0xB4)
VERDE       = RGBColor(0x17, 0x86, 0x4E)
VERDE_CLARO = RGBColor(0xE2, 0xEF, 0xDA)
ROJO        = RGBColor(0xC0, 0x39, 0x2B)
ROJO_CLARO  = RGBColor(0xFC, 0xE4, 0xD6)
GRIS_CLARO  = RGBColor(0xF2, 0xF2, 0xF2)
GRIS_MEDIO  = RGBColor(0xD9, 0xD9, 0xD9)
GRIS_OSCURO = RGBColor(0x66, 0x66, 0x66)
BLANCO      = RGBColor(0xFF, 0xFF, 0xFF)

# ── Leer datos ───────────────────────────────────────────
raw_matriz = pd.read_excel(archivo, sheet_name="Matriz", header=None)
categorias_row      = raw_matriz.iloc[1, 2:].ffill().tolist()
tecnologias_nombres = raw_matriz.iloc[2, 2:].tolist()
tech_categoria = {t: c for t, c in zip(tecnologias_nombres, categorias_row)}

# Todas las tecnologías usadas en alguna aplicación
matriz = pd.read_excel(archivo, sheet_name="Matriz", header=2, index_col=1)
matriz = matriz.drop(columns=[matriz.columns[0]], errors="ignore").dropna(how="all").fillna(0)

personas = pd.read_excel(archivo, sheet_name="MtzPerson", header=2, index_col=1)
personas = personas.drop(columns=[personas.columns[0]], errors="ignore").dropna(how="all").fillna(0)

tecnologias = matriz.columns.intersection(personas.columns)
matriz   = matriz[tecnologias]
personas = personas[tecnologias]

# Universo de tecnologías relevantes (usadas en al menos 1 app)
techs_relevantes = set(tecnologias[(matriz == 1).any(axis=0)])
techs_no_usadas_banco = set(tecnologias) - techs_relevantes

# ── Helpers ──────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_rgb):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=10, bold=False,
             color=RGBColor(0,0,0), align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def altura_bloque_categoria(n_techs, chips_por_fila=4):
    return 0.28 + 0.22 * max(1, -(-n_techs // chips_por_fila))  # ceil div

def construir_bloques_categoria(techs_set):
    por_cat = {}
    for t in techs_set:
        cat = str(tech_categoria.get(t, "Otros")).strip()
        por_cat.setdefault(cat, []).append(t)
    return [(cat, sorted(vals)) for cat, vals in sorted(por_cat.items(), key=lambda x: x[0])]

def seccion_gaps_compacta(slide, titulo, techs_set, col_x, col_y, col_w, col_h, fondo_rgb, titulo_rgb):
    """Renderiza brechas en formato compacto de texto para asegurar 1 slide por persona."""
    def abreviar_tech(txt):
        t = str(txt).strip()
        reemplazos = {
            " (legacy)": "",
            "OpenJDK": "JDK",
            "Microservicios": "MS",
            "Hardware Security Module": "HSM",
            "Active Directory": "AD",
            "On-Premise": "OnPrem",
            "Balanceador": "Bal.",
            "generación": "gen.",
            "Mensajeria": "Msg",
        }
        for k, v in reemplazos.items():
            t = t.replace(k, v)
        t = t.replace(" / ", "/")
        while "  " in t:
            t = t.replace("  ", " ")
        return t

    def abreviar_categoria(txt):
        c = str(txt).strip()
        c = c.replace("Arquitectura", "Arq.")
        c = c.replace("Infraestructura", "Infra")
        c = c.replace("Seguridad", "Seg.")
        c = c.replace("Plataforma", "Plat.")
        return c[:34]

    add_rect(slide, col_x, col_y, col_w, 0.35, titulo_rgb)
    add_text(slide, titulo, col_x + 0.1, col_y + 0.04, col_w - 0.2, 0.3,
             size=10, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

    body_y = col_y + 0.4
    body_h = max(0.2, col_h - 0.4)
    add_rect(slide, col_x, body_y, col_w, body_h, fondo_rgb)

    tech_groups = {}
    for t in techs_set:
        tech = abreviar_tech(t)
        cat = abreviar_categoria(tech_categoria.get(t, "Otros"))
        tech_groups.setdefault(cat, set()).add(tech)

    if not tech_groups:
        add_text(slide, "✓ Sin brechas en esta área", col_x + 0.12, body_y + 0.08,
                 col_w - 0.24, 0.3, size=9, color=VERDE)
        return

    n = sum(len(v) for v in tech_groups.values())
    n_cols = 2 if n <= 70 else (3 if n <= 120 else 4)
    font_size = 8 if n <= 70 else (7 if n <= 120 else 6)

    # Bloques por categoría para mantener encabezado + tecnologías asociadas
    blocks = []
    for cat in sorted(tech_groups.keys(), key=lambda x: x.lower()):
        items = sorted(tech_groups[cat], key=lambda x: x.lower())
        blocks.append((cat, items))

    cols = [[] for _ in range(n_cols)]
    weights = [0] * n_cols
    for cat, items in blocks:
        idx = weights.index(min(weights))
        cols[idx].append((cat, items))
        weights[idx] += 1 + len(items)

    gap = 0.10
    col_inner_w = (col_w - 0.2 - (gap * (n_cols - 1))) / n_cols
    start_x = col_x + 0.1

    for i, blocks_in_col in enumerate(cols):
        cx = start_x + i * (col_inner_w + gap)
        txb = slide.shapes.add_textbox(Inches(cx), Inches(body_y + 0.04),
                                       Inches(col_inner_w), Inches(body_h - 0.08))
        txb.word_wrap = True
        tf = txb.text_frame
        tf.word_wrap = True
        tf.clear()

        first = True
        for cat, items in blocks_in_col:
            p_cat = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p_cat.alignment = PP_ALIGN.LEFT
            run_cat = p_cat.add_run()
            run_cat.text = cat[:34]
            run_cat.font.size = Pt(min(10, font_size + 1))
            run_cat.font.bold = True
            run_cat.font.color.rgb = RGBColor(0x4A, 0x2A, 0x2A)

            for tech in items:
                p_item = tf.add_paragraph()
                p_item.alignment = PP_ALIGN.LEFT
                run_item = p_item.add_run()
                run_item.text = f"• {tech[:32]}"
                run_item.font.size = Pt(font_size)
                run_item.font.bold = False
                run_item.font.color.rgb = RGBColor(0x2A, 0x2A, 0x2A)

def seccion_categorias(slide, titulo, techs_set, col_x, col_y, col_w, fondo_rgb, titulo_rgb,
                      cat_blocks=None, chips_por_fila=4, empty_text="✓ Sin brechas en esta área"):
    """Dibuja un bloque con título y tecnologías agrupadas por categoría."""
    # Título del bloque
    add_rect(slide, col_x, col_y, col_w, 0.35, titulo_rgb)
    add_text(slide, titulo, col_x + 0.1, col_y + 0.04, col_w - 0.2, 0.3,
             size=10, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

    y = col_y + 0.4

    if cat_blocks is None:
        cat_blocks = construir_bloques_categoria(techs_set)

    if not cat_blocks:
        add_rect(slide, col_x, y, col_w, 0.4, fondo_rgb)
        add_text(slide, empty_text, col_x + 0.1, y + 0.08,
                 col_w - 0.2, 0.3, size=9, color=VERDE)
        return y + 0.48

    for cat, techs_list in cat_blocks:
        n_techs = len(techs_list)
        altura_bloque = altura_bloque_categoria(n_techs, chips_por_fila)

        add_rect(slide, col_x, y, col_w, altura_bloque, fondo_rgb)

        # Nombre categoría
        add_text(slide, cat, col_x + 0.1, y + 0.02, col_w - 0.2, 0.22,
                 size=8, bold=True, color=titulo_rgb)

        # Chips de tecnología
        chip_w = (col_w - 0.2) / chips_por_fila - 0.05
        for i, tech in enumerate(techs_list):
            cx = col_x + 0.1 + (i % chips_por_fila) * (chip_w + 0.05)
            cy = y + 0.24 + (i // chips_por_fila) * 0.22
            add_rect(slide, cx, cy, chip_w, 0.18, BLANCO)
            add_text(slide, str(tech).strip()[:22], cx + 0.03, cy + 0.01, chip_w - 0.06, 0.17,
                     size=7, color=RGBColor(0x30, 0x30, 0x30))

        y += altura_bloque + 0.08

    return y

# ── Presentación ─────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# Portada
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 7.5, AZUL_OSCURO)
add_text(slide, "Path de Capacitación", 1, 2.0, 11, 1.2,
         size=40, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
add_text(slide, "Conocimientos actuales y brechas tecnológicas por persona",
         1, 3.3, 11, 0.8, size=18, color=RGBColor(0xBD,0xD7,0xEE), align=PP_ALIGN.CENTER)
add_text(slide, "Banco Santander Chile  ·  2026",
         1, 4.3, 11, 0.5, size=14, color=RGBColor(0xBD,0xD7,0xEE), align=PP_ALIGN.CENTER)

# Slide por persona
for persona_nombre, persona_row in personas.iterrows():
    techs_conocidas = set(tecnologias[persona_row == 1]) & techs_relevantes
    techs_gaps = techs_relevantes - techs_conocidas
    techs_no_usadas_persona = set(tecnologias[persona_row == 1]) & techs_no_usadas_banco

    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 7.5, GRIS_CLARO)
    add_rect(slide, 0, 0, 13.33, 1.0, AZUL_OSCURO)

    # Encabezado
    add_text(slide, str(persona_nombre), 0.3, 0.08, 9, 0.55,
             size=22, bold=True, color=BLANCO)
    add_text(slide, f"{len(techs_conocidas)} conocidas  |  {len(techs_gaps)} por aprender  |  {len(techs_no_usadas_persona)} no usadas en Banco  |  {len(techs_relevantes)} total",
             0.3, 0.62, 12, 0.3, size=11, color=RGBColor(0xBD,0xD7,0xEE))

    # Leyenda visual de colores (encabezado)
    add_text(slide, "Leyenda:", 9.0, 0.08, 1.0, 0.2,
             size=9, bold=True, color=BLANCO)
    add_rect(slide, 9.9, 0.09, 0.18, 0.12, ROJO)
    add_text(slide, "Brecha por aprender", 10.12, 0.075, 1.9, 0.2,
             size=8, color=BLANCO)
    add_rect(slide, 9.9, 0.30, 0.18, 0.12, GRIS_MEDIO)
    add_text(slide, "Conocida, no usada en Banco", 10.12, 0.285, 3.0, 0.2,
             size=8, color=BLANCO)

    # Barra resumen global
    pct_global = len(techs_conocidas) / len(techs_relevantes) * 100 if techs_relevantes else 0
    add_rect(slide, 0, 0.95, 13.33, 0.08, GRIS_MEDIO)
    add_rect(slide, 0, 0.95, 13.33 * pct_global / 100, 0.08, VERDE)

    # Columna izquierda: conocidas + no identificadas en banco
    y_left = seccion_categorias(slide, f"✓  TECNOLOGÍAS QUE YA CONOCE  ({len(techs_conocidas)})",
                                techs_conocidas, 0.2, 1.1, 6.4, VERDE_CLARO, VERDE)

    if techs_no_usadas_persona:
        seccion_categorias(
            slide,
            f"ℹ  TECNOLOGÍA CONOCIDA PERO NO IDENTIFICADA EN BANCO  ({len(techs_no_usadas_persona)})",
            techs_no_usadas_persona, 0.2, y_left + 0.02, 6.4, GRIS_MEDIO, GRIS_OSCURO
        )

    # Columna derecha: gaps en formato compacto (sin paginación)
    seccion_gaps_compacta(slide, f"✗  TECNOLOGÍAS A APRENDER  ({len(techs_gaps)})",
                          techs_gaps, 6.8, 1.1, 6.3, 6.1, ROJO_CLARO, ROJO)

output_file = "PathCapacitacion.pptx"
try:
    prs.save(output_file)
except PermissionError:
    output_file = "PathCapacitacion_actualizado.pptx"
    prs.save(output_file)

print(f"Presentación generada: {output_file}  ({len(personas)} personas)")